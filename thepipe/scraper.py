from typing import Any, Callable, Dict, Iterable, List, Optional, Tuple, Union, cast
import base64
from concurrent.futures import ThreadPoolExecutor, as_completed
from collections import OrderedDict
from io import BytesIO, StringIO
import math
import re
import fnmatch
import os
import tempfile
from urllib.parse import urlparse
import zipfile
from PIL import Image
import requests
import json
from .core import (
    HOST_IMAGES,
    Chunk,
    make_image_url,
    DEFAULT_AI_MODEL,
)
from .chunker import (
    chunk_by_page,
    chunk_by_document,
    chunk_by_section,
    chunk_semantic,
    chunk_by_keywords,
    chunk_by_length,
    chunk_agentic,
)
import tempfile
import mimetypes
import dotenv
from magika import Magika
import markdownify
import fitz
from openai import OpenAI
from openai.types.chat.chat_completion_message_param import ChatCompletionMessageParam

dotenv.load_dotenv()

FOLDERS_TO_IGNORE = {
    "*node_modules*",
    "*.git*",
    "*venv*",
    "*.vscode*",
    "*pycache*",
    "*.ipynb_checkpoints",
}
FILES_TO_IGNORE = {
    ".gitignore",
    "*.bin",
    # Python compiled files
    "*.pyc",
    "*.pyo",
    "*.pyd",
    # Shared libraries and binaries
    "*.so",
    "*.dll",
    "*.exe",
    # Archives and packages
    "*.tar",
    "*.tar.gz",
    "*.egg-info",
    "package-lock.json",
    "package.json",
    # Lock, log, and metadata files
    "*.lock",
    "*.log",
    "Pipfile.lock",
    "requirements.lock",
    "*.exe",
    "*.dll",
    ".DS_Store",
    "Thumbs.db",
}
GITHUB_TOKEN: Optional[str] = os.getenv("GITHUB_TOKEN", None)
USER_AGENT_STRING: str = os.getenv(
    "USER_AGENT_STRING",
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/58.0.3029.110 Safari/537.3",
)
MAX_WHISPER_DURATION = int(os.getenv("MAX_WHISPER_DURATION", 600))  # 10 minutes

TWITTER_DOMAINS = {
    "https://twitter.com",
    "https://www.twitter.com",
    "https://x.com",
    "https://www.x.com",
}
YOUTUBE_DOMAINS = {"https://www.youtube.com", "https://youtube.com"}
GITHUB_DOMAINS = {"https://github.com", "https://www.github.com"}
SCRAPING_PROMPT = os.getenv(
    "SCRAPING_PROMPT",
    """A document is given. Please output the entire extracted contents from the document in detailed markdown format.
Your accuracy is very important. Please be careful to not miss any content from the document.
Be sure to correctly output a comprehensive format markdown for all the document contents (including, but not limited to, headers, paragraphs, lists, tables, menus, equations, full text contents, titles, subtitles, appendices, page breaks, columns, footers, page numbers, watermarks, footnotes, captions, annotations, images, figures, charts, shapes, form fields, content controls, signatures, etc.)
Always reply immediately with only markdown.
Do not give the markdown in a code block. Simply output the raw markdown immediately.
Do not output anything else.""",
)
FILESIZE_LIMIT_MB = int(os.getenv("FILESIZE_LIMIT_MB", 50))  # for url scraping only


def _load_whisper():
    try:
        import whisper
    except ImportError as exc:  # pragma: no cover - optional dependency
        raise ImportError(
            "Audio and video transcription requires the optional dependency `openai-whisper`. "
            "Install it with `pip install thepipe-api[audio]` or include the `gpu` extra."
        ) from exc

    return whisper


def detect_source_mimetype(source: str) -> str:
    # try to detect the file type by its extension
    _, extension = os.path.splitext(source)
    if extension:
        if extension == ".ipynb":
            # special case for notebooks, mimetypes is not familiar
            return "application/x-ipynb+json"
        guessed_mimetype, _ = mimetypes.guess_type(source)
        if guessed_mimetype:
            return guessed_mimetype
    # if that fails, try AI detection with Magika
    magika = Magika()
    with open(source, "rb") as file:
        result = magika.identify_bytes(file.read())
    mimetype = result.output.mime_type
    return mimetype


def scrape_file(
    filepath: str,
    verbose: bool = False,
    chunking_method: Optional[Callable[[List[Chunk]], List[Chunk]]] = chunk_by_page,
    openai_client: Optional[OpenAI] = None,
    model: str = DEFAULT_AI_MODEL,
    include_input_images: bool = True,
    include_output_images: bool = True,
) -> List[Chunk]:
    """
    Scrapes a file and returns a list of Chunk objects containing the text and images extracted from the file.

    Parameters
    ----------
    filepath : str
        The path to the file to scrape.
    verbose : bool, optional
        If ``True``, prints verbose output.
    chunking_method : Callable, optional
        A function to chunk the scraped content. Defaults to chunk_by_page.
    openai_client : OpenAI, optional
        An OpenAI client instance for LLM processing. If provided, uses VLM to scrape PDFs.
    model : str, optional
        The LLM model name to use for processing. Defaults to DEFAULT_AI_MODEL.
    include_input_images : bool, optional
        If ``True``, includes input images in the messages sent to the LLM.
    include_output_images : bool, optional
        If ``True``, includes output images in the returned chunks.
    Returns
    -------
    List[Chunk]
        A list of Chunk objects containing the scraped content.
    """
    # returns chunks of scraped content from any source (file, URL, etc.)
    scraped_chunks = []
    source_mimetype = detect_source_mimetype(filepath)
    if source_mimetype is None:
        if verbose:
            print(f"[thepipe] Unsupported source type: {filepath}")
        return scraped_chunks
    if verbose:
        print(f"[thepipe] Scraping {source_mimetype}: {filepath}...")
    if source_mimetype == "application/pdf":
        scraped_chunks = scrape_pdf(
            file_path=filepath,
            verbose=verbose,
            model=model,
            openai_client=openai_client,
            include_input_images=include_input_images,
            include_output_images=include_output_images,
        )
    elif (
        source_mimetype
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        scraped_chunks = scrape_docx(
            file_path=filepath,
            verbose=verbose,
            include_output_images=include_output_images,
        )
    elif (
        source_mimetype
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        scraped_chunks = scrape_pptx(
            file_path=filepath,
            verbose=verbose,
            include_output_images=include_output_images,
        )
    elif source_mimetype.startswith("image/"):
        scraped_chunks = scrape_image(file_path=filepath)
    elif (
        source_mimetype.startswith("application/vnd.ms-excel")
        or source_mimetype
        == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    ):
        scraped_chunks = scrape_spreadsheet(
            file_path=filepath, source_type=source_mimetype
        )
    elif source_mimetype == "application/x-ipynb+json":
        scraped_chunks = scrape_ipynb(
            file_path=filepath,
            verbose=verbose,
            include_output_images=include_output_images,
        )
    elif (
        source_mimetype == "application/zip"
        or source_mimetype == "application/x-zip-compressed"
    ):
        scraped_chunks = scrape_zip(
            file_path=filepath,
            verbose=verbose,
            openai_client=openai_client,
            include_input_images=include_input_images,
            include_output_images=include_output_images,
        )
    elif source_mimetype.startswith("video/"):
        scraped_chunks = scrape_video(
            file_path=filepath,
            verbose=verbose,
            include_output_images=include_output_images,
        )
    elif source_mimetype.startswith("audio/"):
        scraped_chunks = scrape_audio(file_path=filepath, verbose=verbose)
    elif source_mimetype.startswith("text/html"):
        scraped_chunks = scrape_html(
            file_path=filepath,
            verbose=verbose,
            include_output_images=include_output_images,
        )
    elif source_mimetype.startswith("text/"):
        scraped_chunks = scrape_plaintext(file_path=filepath)
    else:
        try:
            scraped_chunks = scrape_plaintext(file_path=filepath)
        except Exception as e:
            if verbose:
                print(f"[thepipe] Error extracting from {filepath}: {e}")
    if verbose:
        if scraped_chunks:
            print(f"[thepipe] Extracted from {filepath}")
        else:
            print(f"[thepipe] No content extracted from {filepath}")
    if chunking_method:
        scraped_chunks = chunking_method(scraped_chunks)
    return scraped_chunks


def scrape_html(
    file_path: str,
    verbose: bool = False,
    include_output_images: bool = True,
) -> List[Chunk]:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
        html_content = file.read()
    markdown_content = markdownify.markdownify(html_content, heading_style="ATX")
    images = get_images_from_markdown(html_content) if include_output_images else []
    return [Chunk(path=file_path, text=markdown_content, images=images)]


def scrape_plaintext(file_path: str) -> List[Chunk]:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
        text = file.read()
    return [Chunk(path=file_path, text=text)]


def scrape_directory(
    dir_path: str,
    inclusion_pattern: Optional[str] = None,
    verbose: bool = False,
    openai_client: Optional[OpenAI] = None,
    model: str = DEFAULT_AI_MODEL,
    include_input_images: bool = True,
    include_output_images: bool = True,
) -> List[Chunk]:
    """
    inclusion_pattern: Optional regex string; only files whose path matches this pattern will be scraped.
    By default, ignores all files in baked-in constants FOLDERS_TO_IGNORE and FILES_TO_IGNORE.
    """
    # compile the include pattern once
    pattern = re.compile(inclusion_pattern) if inclusion_pattern else None
    extraction: List[Chunk] = []

    try:
        for entry in os.scandir(dir_path):
            path = entry.path

            # skip ignored directories
            if entry.is_dir() and any(
                fnmatch.fnmatch(entry.name, pat) for pat in FOLDERS_TO_IGNORE
            ):
                if verbose:
                    print(f"[thepipe] Skipping ignored directory: {path}")
                continue

            # skip ignored files
            if entry.is_file() and any(
                fnmatch.fnmatch(entry.name, pat) for pat in FILES_TO_IGNORE
            ):
                if verbose:
                    print(f"[thepipe] Skipping ignored file: {path}")
                continue

            if entry.is_file():
                # if include_pattern is set, skip files that don't match
                if pattern and not pattern.search(path):
                    if verbose:
                        print(f"[thepipe] Skipping non-matching file: {path}")
                    continue

                if verbose:
                    print(f"[thepipe] Scraping file: {path}")
                extraction += scrape_file(
                    filepath=path,
                    verbose=verbose,
                    openai_client=openai_client,
                    model=model,
                    include_input_images=include_input_images,
                    include_output_images=include_output_images,
                )

            elif entry.is_dir():
                # recurse into subdirectory
                if verbose:
                    print(f"[thepipe] Entering directory: {path}")
                extraction += scrape_directory(
                    dir_path=path,
                    inclusion_pattern=inclusion_pattern,
                    verbose=verbose,
                    openai_client=openai_client,
                    model=model,
                    include_input_images=include_input_images,
                    include_output_images=include_output_images,
                )
    except PermissionError as e:
        if verbose:
            print(f"[thepipe] Skipping {dir_path} (permission denied): {e}")

    return extraction


def scrape_zip(
    file_path: str,
    inclusion_pattern: Optional[str] = None,
    verbose: bool = False,
    openai_client: Optional[OpenAI] = None,
    include_input_images: bool = True,
    include_output_images: bool = True,
) -> List[Chunk]:
    chunks = []
    with tempfile.TemporaryDirectory() as temp_dir:
        with zipfile.ZipFile(file_path, "r") as zip_ref:
            zip_ref.extractall(temp_dir)
        chunks = scrape_directory(
            dir_path=temp_dir,
            inclusion_pattern=inclusion_pattern,
            verbose=verbose,
            openai_client=openai_client,
            include_input_images=include_input_images,
            include_output_images=include_output_images,
        )
    return chunks


def scrape_pdf(
    file_path: str,
    openai_client: Optional[OpenAI] = None,
    model: str = DEFAULT_AI_MODEL,
    verbose: Optional[bool] = False,
    include_input_images: bool = True,
    include_output_images: bool = True,
    image_scale: float = 1.0,
) -> List[Chunk]:
    chunks: List[Chunk] = []

    # Branch 1 – VLM path (OpenAI client supplied)
    if openai_client is not None:
        with open(file_path, "rb") as fp:
            pdf_bytes = fp.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        num_pages = len(doc)

        if verbose:
            print(
                f"[thepipe] Scraping PDF: {file_path} "
                f"({num_pages} pages) with model {model}"
            )

        # Inner worker – processes one page
        def _process_page(page_num: int) -> Tuple[int, str, Optional[Image.Image]]:
            page = doc[page_num]
            text = page.get_text()  # type: ignore[attr-defined]

            # Build message for the LLM
            msg_content: List[Dict[str, Union[Dict[str, str], str]]] = [
                {
                    "type": "text",
                    "text": f"```\n{text}\n```\n{SCRAPING_PROMPT}",
                }
            ]

            image: Optional[Image.Image] = None
            if include_input_images or include_output_images:
                mat = fitz.Matrix(image_scale, image_scale)
                pix = page.get_pixmap(matrix=mat, alpha=False)  # type: ignore[attr-defined]  # noqa: E501
                image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

                if include_input_images:
                    encoded = make_image_url(image, host_images=HOST_IMAGES)
                    msg_content.append(
                        {
                            "type": "image_url",
                            "image_url": {"url": encoded, "detail": "high"},
                        }
                    )

            messages = cast(
                Iterable[ChatCompletionMessageParam],
                [{"role": "user", "content": msg_content}],
            )

            response = openai_client.chat.completions.create(
                model=model, messages=messages
            )

            llm_response = response.choices[0].message.content
            if not llm_response:
                raise RuntimeError("Empty LLM response.")

            llm_response = llm_response.strip()
            if llm_response.startswith("```markdown"):
                llm_response = llm_response[len("```markdown") :]
            elif llm_response.startswith("```"):
                llm_response = llm_response[len("```") :]
            if llm_response.endswith("```"):
                llm_response = llm_response[: -len("```")]

            return (
                page_num,
                llm_response,
                image if include_output_images else None,
            )

        # Parallel extraction
        max_workers = (os.cpu_count() or 1) * 2
        if verbose:
            print(f"[thepipe] Using {max_workers} threads for PDF extraction")

        page_results: OrderedDict[int, Tuple[str, Optional[Image.Image]]] = (
            OrderedDict()
        )
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = [executor.submit(_process_page, p) for p in range(num_pages)]
            for fut in as_completed(futures):
                pg, txt, img = fut.result()
                page_results[pg] = (txt, img)

        for pg in sorted(page_results):
            txt, img = page_results[pg]
            chunks.append(Chunk(path=file_path, text=txt, images=[img] if img else []))

        return chunks

    # Branch 2 – no OpenAI client – text-only offline mode
    from pymupdf4llm.helpers.pymupdf_rag import to_markdown  # local import

    doc = fitz.open(file_path)
    md_pages = cast(List[Dict[str, Any]], to_markdown(file_path, page_chunks=True))

    for i in range(doc.page_count):
        text = re.sub(r"\n{3,}", "\n\n", md_pages[i]["text"]).strip()

        images: List[Image.Image] = []
        if include_output_images:
            mat = fitz.Matrix(image_scale, image_scale)
            pix = doc[i].get_pixmap(matrix=mat, alpha=False)  # type: ignore[attr-defined]  # noqa: E501
            images.append(Image.frombytes("RGB", [pix.width, pix.height], pix.samples))

        chunks.append(Chunk(path=file_path, text=text, images=images))

    doc.close()
    return chunks


def get_images_from_markdown(text: str) -> List[Image.Image]:
    image_urls = re.findall(r"!\[.*?\]\((.*?)\)", text)
    images = []
    for url in image_urls:
        extension = os.path.splitext(urlparse(url).path)[1]
        if extension not in {".jpg", ".jpeg", ".png"}:
            # ignore incompatible image extractions
            continue

        try:
            response = requests.get(
                url,
                timeout=10,
                headers={"User-Agent": USER_AGENT_STRING},
            )
            response.raise_for_status()
        except Exception:
            continue

        img = Image.open(BytesIO(response.content))
        images.append(img)
    return images


def scrape_image(file_path: str) -> List[Chunk]:
    img = Image.open(file_path)
    img.load()  # needed to close the file
    chunk = Chunk(path=file_path, images=[img])
    return [chunk]


def scrape_spreadsheet(file_path: str, source_type: str) -> List[Chunk]:
    import pandas as pd

    if source_type == "application/vnd.ms-excel":
        df = pd.read_csv(file_path)
    elif (
        source_type
        == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    ):
        df = pd.read_excel(file_path)
    else:
        raise ValueError("Unsupported file format")
    dicts = df.to_dict(orient="records")
    chunks = []
    for i, item in enumerate(dicts):
        # format each row as json along with the row index
        item["row index"] = i
        item_json = json.dumps(item, indent=4)
        chunks.append(Chunk(path=file_path, text=item_json))
    return chunks


def parse_webpage_with_vlm(
    url: str,
    model: str = DEFAULT_AI_MODEL,
    verbose: Optional[bool] = False,
    openai_client: Optional[OpenAI] = None,
    include_output_images: bool = True,
) -> Chunk:
    if openai_client is None:
        raise ValueError("parse_webpage_with_vlm requires an openai_client argument.")
    from playwright.sync_api import sync_playwright

    with sync_playwright() as p:
        browser = p.chromium.launch()
        context = browser.new_context(user_agent=USER_AGENT_STRING)
        page = context.new_page()
        page.goto(url, wait_until="domcontentloaded")
        if not page.viewport_size:
            page.set_viewport_size({"width": 800, "height": 600})
        if not page.viewport_size:
            raise ValueError(
                "Failed to set viewport size after finding no viewport size"
            )
        viewport_height = page.viewport_size.get("height", 800)
        total_height = page.evaluate("document.body.scrollHeight")
        current_scroll_position = 0
        scrolldowns, max_scrolldowns = 0, 3
        images: List[Image.Image] = []

        while current_scroll_position < total_height and scrolldowns < max_scrolldowns:
            page.wait_for_timeout(200)  # wait for content to load
            screenshot = page.screenshot(full_page=False)
            img = Image.open(BytesIO(screenshot))
            images.append(img)

            current_scroll_position += viewport_height
            page.evaluate(f"window.scrollTo(0, {current_scroll_position})")
            scrolldowns += 1
            total_height = page.evaluate("document.body.scrollHeight")
            if verbose:
                print(
                    f"[thepipe] Scrolled to {current_scroll_position} of {total_height}. Waiting for content to load..."
                )

        browser.close()

    if images:
        # Vertically stack the images
        total_height = sum(img.height for img in images)
        max_width = max(img.width for img in images)
        stacked_image = Image.new("RGB", (max_width, total_height))
        y_offset = 0
        for img in images:
            stacked_image.paste(img, (0, y_offset))
            y_offset += img.height

        # Process the stacked image with VLM
        messages = [
            {
                "role": "user",
                "content": [
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": make_image_url(
                                stacked_image, host_images=HOST_IMAGES
                            ),
                            "detail": "high",
                        },
                    },
                    {"type": "text", "text": SCRAPING_PROMPT},
                ],
            },
        ]
        response = openai_client.chat.completions.create(
            model=model,
            messages=cast(Iterable[ChatCompletionMessageParam], messages),
        )
        llm_response = response.choices[0].message.content
        if not llm_response:
            raise Exception(
                f"Failed to receive a message content from LLM Response: {response}"
            )
        if verbose:
            print(f"[thepipe] LLM response: {llm_response}")
        chunk = Chunk(
            path=url,
            text=llm_response,
            images=[stacked_image] if include_output_images else [],
        )
    else:
        raise ValueError("Model received 0 images from webpage")

    return chunk


def extract_page_content(
    url: str, verbose: bool = False, include_output_images: bool = True
) -> Chunk:
    from bs4 import BeautifulSoup
    from playwright.sync_api import sync_playwright
    import base64
    import requests

    texts: List[str] = []
    images: List[Image.Image] = []

    with sync_playwright() as p:
        browser = p.chromium.launch()
        context = browser.new_context(user_agent=USER_AGENT_STRING)
        page = context.new_page()

        try:
            page.goto(url, wait_until="domcontentloaded", timeout=10000)

            # Wait for content to load
            page.wait_for_timeout(1000)

            # Scroll to load dynamic content
            if not page.viewport_size:
                page.set_viewport_size({"width": 1200, "height": 800})

            viewport_height = page.viewport_size["height"]
            total_height = page.evaluate("document.body.scrollHeight")
            current_scroll_position = 0
            scrolldowns, max_scrolldowns = 0, 5

            while (
                current_scroll_position < total_height and scrolldowns < max_scrolldowns
            ):
                page.wait_for_timeout(500)
                current_scroll_position += viewport_height
                page.evaluate(f"window.scrollTo(0, {current_scroll_position})")
                scrolldowns += 1
                new_height = page.evaluate("document.body.scrollHeight")
                if new_height == total_height:
                    break
                total_height = new_height

            # Extract HTML content
            html_content = page.content()

            # Parse with BeautifulSoup and clean up
            soup = BeautifulSoup(html_content, "html.parser")

            # Remove script and style elements
            for script in soup(["script", "style", "nav", "footer", "header"]):
                script.decompose()

            # Convert to markdown
            markdown_content = markdownify.markdownify(str(soup), heading_style="ATX")

            # Clean up markdown
            markdown_content = re.sub(r"\n{3,}", "\n\n", markdown_content)
            markdown_content = markdown_content.strip()

            if verbose:
                print(
                    f"[thepipe] Extracted {len(markdown_content)} characters from {url}"
                )

            texts.append(markdown_content)

            # Extract images from the page using heuristics
            if include_output_images:
                for img in page.query_selector_all("img"):
                    img_path = img.get_attribute("src")
                    if not img_path:
                        continue
                    if img_path.startswith("data:image"):
                        # Save base64 image to PIL Image
                        try:
                            decoded_data = base64.b64decode(img_path.split(",")[1])
                            image = Image.open(BytesIO(decoded_data))
                            images.append(image)
                        except Exception as e:
                            if verbose:
                                print(
                                    f"[thepipe] Ignoring error loading base64 image: {e}"
                                )
                            continue
                    else:
                        try:
                            # Try direct URL first
                            response = requests.get(
                                img_path,
                                timeout=10,
                                headers={"User-Agent": USER_AGENT_STRING},
                            )
                            response.raise_for_status()
                            image = Image.open(BytesIO(response.content))
                            images.append(image)
                        except Exception as e:
                            if verbose:
                                print(f"[thepipe] Error loading image {img_path}: {e}")
                                print("[thepipe] Attempting to load path with schema.")

                            # Try with schema if path is relative
                            if not img_path.startswith(("http://", "https://")):
                                try:
                                    # Remove leading slashes
                                    while img_path.startswith("/"):
                                        img_path = img_path[1:]

                                    # Try with just the scheme
                                    parsed_url = urlparse(url)
                                    path_with_schema = (
                                        f"{parsed_url.scheme}://{img_path}"
                                    )
                                    response = requests.get(
                                        path_with_schema,
                                        timeout=10,
                                        headers={"User-Agent": USER_AGENT_STRING},
                                    )
                                    response.raise_for_status()
                                    image = Image.open(BytesIO(response.content))
                                    images.append(image)
                                except Exception as e:
                                    if verbose:
                                        print(
                                            f"[thepipe] Error loading image {img_path} with schema: {e}"
                                        )
                                        print(
                                            "[thepipe] Attempting to load with schema and netloc."
                                        )

                                    try:
                                        # Try with scheme and netloc
                                        path_with_schema_and_netloc = f"{parsed_url.scheme}://{parsed_url.netloc}/{img_path}"
                                        response = requests.get(
                                            path_with_schema_and_netloc,
                                            timeout=10,
                                            headers={"User-Agent": USER_AGENT_STRING},
                                        )
                                        response.raise_for_status()
                                        image = Image.open(BytesIO(response.content))
                                        images.append(image)
                                    except Exception as e:
                                        if verbose:
                                            print(
                                                f"[thepipe] Final attempt failed for image {img_path}: {e}"
                                            )
                                        continue
                            else:
                                if verbose:
                                    print(
                                        f"[thepipe] Skipping image {img_path} - all attempts failed"
                                    )
                                continue

        except Exception as e:
            if verbose:
                print(f"[thepipe] Error scraping {url}: {e}")
            # Fallback to simple requests
            try:
                response = requests.get(
                    url, headers={"User-Agent": USER_AGENT_STRING}, timeout=30
                )
                response.raise_for_status()
                soup = BeautifulSoup(response.content, "html.parser")

                # Remove unwanted elements
                for script in soup(["script", "style", "nav", "footer", "header"]):
                    script.decompose()

                markdown_content = markdownify.markdownify(
                    str(soup), heading_style="ATX"
                )
                markdown_content = re.sub(r"\n{3,}", "\n\n", markdown_content).strip()
                texts.append(markdown_content)

                if verbose:
                    print(
                        f"[thepipe] Fallback extraction got {len(markdown_content)} characters"
                    )

            except Exception as fallback_e:
                if verbose:
                    print(f"[thepipe] Fallback also failed: {fallback_e}")
                texts.append("")

        finally:
            browser.close()

    text = "\n".join(texts).strip()
    return Chunk(path=url, text=text, images=images)


def scrape_url(
    url: str,
    verbose: bool = False,
    chunking_method: Callable[[List[Chunk]], List[Chunk]] = chunk_by_page,
    openai_client: Optional[OpenAI] = None,
    model: str = DEFAULT_AI_MODEL,
    include_input_images: bool = True,
    include_output_images: bool = True,
) -> List[Chunk]:
    if any(url.startswith(domain) for domain in TWITTER_DOMAINS):
        extraction = scrape_tweet(url=url, include_output_images=include_output_images)
        return extraction
    elif any(url.startswith(domain) for domain in YOUTUBE_DOMAINS):
        extraction = scrape_youtube(
            youtube_url=url,
            verbose=verbose,
            include_output_images=include_output_images,
        )
        return extraction
    elif any(url.startswith(domain) for domain in GITHUB_DOMAINS):
        extraction = scrape_github(
            github_url=url,
            verbose=verbose,
            openai_client=openai_client,
            model=model,
            include_input_images=include_input_images,
            include_output_images=include_output_images,
        )
        return extraction
    _, extension = os.path.splitext(urlparse(url).path)
    if extension and extension not in {".html", ".htm", ".php", ".asp", ".aspx"}:
        # if url leads to a file, attempt to download it and scrape it
        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = os.path.join(temp_dir, os.path.basename(url))
            response = requests.get(url)
            # verify the ingress/egress with be within limits, if there are any set
            response_length = int(response.headers.get("Content-Length", 0))
            if FILESIZE_LIMIT_MB and response_length > FILESIZE_LIMIT_MB * 1024 * 1024:
                raise ValueError(f"File size exceeds {FILESIZE_LIMIT_MB} MB limit.")
            with open(file_path, "wb") as file:
                file.write(response.content)
            chunks = scrape_file(
                filepath=file_path,
                verbose=verbose,
                chunking_method=chunking_method,
                openai_client=openai_client,
                model=model,
                include_input_images=include_input_images,
                include_output_images=include_output_images,
            )
        return chunks
    else:
        # if url leads to web content, scrape it directly
        if openai_client and include_input_images:
            chunk = parse_webpage_with_vlm(
                url=url,
                verbose=verbose,
                model=model,
                openai_client=openai_client,
                include_output_images=include_output_images,
            )
        else:
            chunk = extract_page_content(
                url=url, verbose=verbose, include_output_images=include_output_images
            )
        chunks = chunking_method([chunk])
        # if no text or images were extracted, return error
        if not any(chunk.text for chunk in chunks) and not any(
            chunk.images for chunk in chunks
        ):
            raise ValueError("No content extracted from URL.")
        return chunks


def format_timestamp(seconds, chunk_index, chunk_duration):
    # helper function to format the timestamp.
    total_seconds = chunk_index * chunk_duration + seconds
    hours = int(total_seconds // 3600)
    minutes = int((total_seconds % 3600) // 60)
    seconds = total_seconds % 60
    milliseconds = int((seconds - int(seconds)) * 1000)
    return f"{hours:02}:{minutes:02}:{int(seconds):02}.{milliseconds:03}"


def scrape_video(
    file_path: str,
    verbose: bool = False,
    include_output_images: bool = True,
) -> List[Chunk]:
    whisper = _load_whisper()
    from moviepy.editor import VideoFileClip

    # Splits the video into chunks of length MAX_WHISPER_DURATION, extracts
    # one representative frame from the start of each chunk, and then transcribes
    # that chunk.
    model = whisper.load_model("base")
    video = VideoFileClip(file_path)
    num_chunks = math.ceil(video.duration / MAX_WHISPER_DURATION)
    chunks = []

    try:
        for i in range(num_chunks):
            # Calculate the start and end time of the chunk
            start_time = i * MAX_WHISPER_DURATION
            end_time = start_time + MAX_WHISPER_DURATION
            if end_time > video.duration:
                end_time = video.duration

            # Extract a frame from the start of the chunk
            image = None
            if include_output_images:
                frame = video.get_frame(start_time)
                image = Image.fromarray(frame)

            # Save the audio to a temporary .wav file
            with tempfile.NamedTemporaryFile(
                suffix=".wav", delete=False
            ) as temp_audio_file:
                audio_path = temp_audio_file.name

            audio = video.subclip(start_time, end_time).audio  # type: ignore[attr-defined]
            transcription = None

            if audio is not None:
                audio.write_audiofile(audio_path, codec="pcm_s16le")
                result = model.transcribe(audio=audio_path, verbose=verbose)

                # Format transcription with timestamps
                formatted_transcription = []
                for segment in cast(List[Dict[str, Any]], result["segments"]):
                    seg_start = format_timestamp(
                        segment["start"], i, MAX_WHISPER_DURATION
                    )
                    seg_end = format_timestamp(segment["end"], i, MAX_WHISPER_DURATION)
                    formatted_transcription.append(
                        f"[{seg_start} --> {seg_end}]  {segment['text']}"
                    )

                transcription = "\n".join(formatted_transcription)
                os.remove(audio_path)

            # Only add chunks if there is either text or images
            if transcription or image:
                chunks.append(
                    Chunk(
                        path=file_path,
                        text=transcription if transcription else None,
                        images=[image] if image else [],
                    )
                )
    finally:
        video.close()

    return chunks


def scrape_youtube(
    youtube_url: str,
    verbose: bool = False,
    include_output_images: bool = True,
) -> List[Chunk]:
    from pytube import YouTube

    with tempfile.TemporaryDirectory() as temp_dir:
        filename = "temp_video.mp4"
        yt = YouTube(youtube_url)
        stream = yt.streams.filter(progressive=True, file_extension="mp4").first()
        if stream is None:
            raise ValueError("No progressive stream for video found.")
        stream.download(temp_dir, filename=filename)
        video_path = os.path.join(temp_dir, filename)
        chunks = scrape_video(
            file_path=video_path,
            verbose=verbose,
            include_output_images=include_output_images,
        )
    return chunks


def scrape_audio(file_path: str, verbose: bool = False) -> List[Chunk]:
    whisper = _load_whisper()

    model = whisper.load_model("base")
    result = model.transcribe(audio=file_path, verbose=verbose)
    segments = cast(List[Dict[str, Any]], result.get("segments", []))

    transcript: List[str] = []
    for segment in segments:
        start = format_timestamp(segment["start"], 0, 0)
        end = format_timestamp(segment["end"], 0, 0)
        if segment["text"].strip():
            transcript.append(f"[{start} --> {end}]  {segment['text']}")
    # join the formatted transcription into a single string
    transcript_text = "\n".join(transcript)
    return [Chunk(path=file_path, text=transcript_text)]


def scrape_github(
    github_url: str,
    inclusion_pattern: Optional[str] = None,
    branch: str = "main",
    verbose: bool = False,
    openai_client: Optional[OpenAI] = None,
    model: str = DEFAULT_AI_MODEL,
    include_input_images: bool = True,
    include_output_images: bool = True,
) -> List[Chunk]:
    files_contents: List[Chunk] = []
    if not GITHUB_TOKEN:
        raise ValueError("GITHUB_TOKEN environment variable is not set.")
    # make new tempdir for cloned repo
    with tempfile.TemporaryDirectory() as temp_dir:
        # requires git
        exit_code = os.system(
            f'git clone --branch "{branch}" --single-branch {github_url} "{temp_dir}" --quiet'
        )
        if exit_code != 0:
            raise RuntimeError(
                f"git clone failed for {github_url} at branch '{branch}'. "
                "Verify the repository URL and branch name."
            )
        files_contents = scrape_directory(
            dir_path=temp_dir,
            inclusion_pattern=inclusion_pattern,
            verbose=verbose,
            openai_client=openai_client,
            model=model,
            include_input_images=include_input_images,
            include_output_images=include_output_images,
        )
    return files_contents


def scrape_docx(
    file_path: str,
    verbose: bool = False,
    include_output_images: bool = True,
) -> List[Chunk]:
    from docx import Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table, _Cell
    from docx.text.paragraph import Paragraph
    import csv
    import io

    # helper function to iterate through blocks in the document
    def iter_block_items(parent):
        if parent.__class__.__name__ == "Document":
            parent_elm = parent.element.body
        elif parent.__class__.__name__ == "_Cell":
            parent_elm = parent._tc
        else:
            raise ValueError("Unsupported parent type")
        # iterate through each child element in the parent element
        for child in parent_elm.iterchildren():
            child_elem_class_name = child.__class__.__name__
            if verbose:
                print(f"[thepipe] Found element in docx: {child_elem_class_name}")
            if child_elem_class_name == "CT_P":
                yield Paragraph(child, parent)
            elif child_elem_class_name == "CT_Tbl":
                yield Table(child, parent)

    # helper function to read tables in the document
    def read_docx_tables(tab):
        vf = StringIO()
        writer = csv.writer(vf)
        for row in tab.rows:
            writer.writerow(cell.text for cell in row.cells)
        vf.seek(0)
        return vf.getvalue()

    # read the document
    document = Document(file_path)
    chunks = []
    image_counter = 0

    # Define namespaces
    nsmap = {
        "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
        "pic": "http://schemas.openxmlformats.org/drawingml/2006/picture",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }

    try:
        # scrape each block in the document to create chunks
        # A block can be a paragraph, table, or image
        for block in iter_block_items(document):
            block_texts = []
            block_images = []
            if isinstance(block, Paragraph):
                block_texts.append(block.text)
                # "runs" are the smallest units in a paragraph
                for run in block.runs:
                    if "pic:pic" in run.element.xml and include_output_images:
                        # extract images from the paragraph
                        for pic in run.element.findall(".//pic:pic", nsmap):
                            cNvPr = pic.find(".//pic:cNvPr", nsmap)
                            name_attr = (
                                cNvPr.get("name")
                                if cNvPr is not None
                                else f"image_{image_counter}"
                            )
                            blip = pic.find(".//a:blip", nsmap)
                            if blip is not None:
                                embed_attr = blip.get(
                                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                                )
                                if embed_attr:
                                    image_part = document.part.related_parts[embed_attr]
                                    image_data = BytesIO(image_part._blob)
                                    image = Image.open(image_data)
                                    image.load()
                                    block_images.append(image)
                                    image_counter += 1
            elif isinstance(block, Table):
                table_text = read_docx_tables(block)
                block_texts.append(table_text)
            if block_texts or block_images:
                block_text = "\n".join(block_texts).strip()
                if block_text or block_images:
                    chunks.append(
                        Chunk(path=file_path, text=block_text, images=block_images)
                    )
    except Exception as e:
        raise ValueError(f"Error processing DOCX file {file_path}: {e}")
    return chunks


def scrape_pptx(
    file_path: str,
    verbose: bool = False,
    include_output_images: bool = True,
) -> List[Chunk]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.shapes.picture import Picture
    from pptx.shapes.autoshape import Shape as AutoShape

    prs = Presentation(file_path)
    chunks = []
    # iterate through each slide in the presentation
    for slide in prs.slides:
        slide_texts = []
        slide_images = []
        # iterate through each shape in the slide
        for shape in slide.shapes:
            if shape.has_text_frame:
                auto_shape = cast(AutoShape, shape)
                for paragraph in auto_shape.text_frame.paragraphs:
                    text = paragraph.text
                    if len(slide_texts) == 0:
                        text = "# " + text  # header for first text of a slide
                    slide_texts.append(text)
            # extract images from shapes
            if include_output_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pic = cast(Picture, shape)
                image_data = pic.image.blob
                image = Image.open(BytesIO(image_data))
                slide_images.append(image)
        # add slide to chunks if it has text or images
        if slide_texts or slide_images:
            text = "\n".join(slide_texts).strip()
            if not include_output_images:
                slide_images = []
            chunks.append(Chunk(path=file_path, text=text, images=slide_images))
    # return all chunks
    return chunks


def scrape_ipynb(
    file_path: str,
    verbose: bool = False,
    include_output_images: bool = True,
) -> List[Chunk]:
    with open(file_path, "r", encoding="utf-8") as file:
        notebook = json.load(file)
    chunks = []
    # parse cells in the notebook
    for cell in notebook["cells"]:
        texts = []
        images: List[Image.Image] = []
        cell_type = cell["cell_type"]
        # parse cell content based on type
        if verbose:
            print(f"[thepipe] Scraping cell {cell_type} from {file_path}")
        if cell_type == "markdown":
            text = "".join(cell["source"])
            if include_output_images:
                images = get_images_from_markdown(text)
            texts.append(text)
        elif cell_type == "code":
            source = "".join(cell["source"])
            texts.append(source)
            output_texts = []
            # code cells can have outputs
            if "outputs" in cell:
                for output in cell["outputs"]:
                    if (
                        include_output_images
                        and "data" in output
                        and "image/png" in output["data"]
                    ):
                        image_data = output["data"]["image/png"]
                        image = Image.open(BytesIO(base64.b64decode(image_data)))
                        images.append(image)
                    elif "data" in output and "text/plain" in output["data"]:
                        output_text = "".join(output["data"]["text/plain"])
                        output_texts.append(output_text)
            if output_texts:
                texts.extend(output_texts)
        elif cell_type == "raw":
            text = "".join(cell["source"])
            texts.append(text)
        if texts or images:
            text = "\n".join(texts).strip()
            chunks.append(Chunk(path=file_path, text=text, images=images))
    return chunks


def scrape_tweet(url: str, include_output_images: bool = True) -> List[Chunk]:
    """
    Magic function from https://github.com/vercel/react-tweet/blob/main/packages/react-tweet/src/api/fetch-tweet.ts
    unofficial, could break at any time
    """

    def get_token(id: str) -> str:
        result = (float(id) / 1e15) * math.pi
        base_36_result = ""
        characters = "0123456789abcdefghijklmnopqrstuvwxyz"
        while result > 0:
            remainder = int(result % (6**2))
            base_36_result = characters[remainder] + base_36_result
            result = (result - remainder) // (6**2)
        base_36_result = re.sub(r"(0+|\.)", "", base_36_result)
        return base_36_result

    tweet_id = url.split("status/")[-1].split("?")[0]
    token = get_token(tweet_id)
    tweet_api_url = "https://cdn.syndication.twimg.com/tweet-result"
    params = {"id": tweet_id, "language": "en", "token": token}
    response = requests.get(tweet_api_url, params=params)
    if response.status_code != 200:
        raise ValueError(f"Failed to fetch tweet. Status code: {response.status_code}")
    tweet_data = response.json()
    # Extract tweet text
    tweet_text = tweet_data.get("text", "")
    # Extract images from tweet
    images: List[Image.Image] = []
    if include_output_images and "mediaDetails" in tweet_data:
        for media in tweet_data["mediaDetails"]:
            image_url = media.get("media_url_https")
            if image_url:
                image_response = requests.get(image_url)
                img = Image.open(BytesIO(image_response.content))
                images.append(img)
    # Create chunks for text and images
    chunk = Chunk(path=url, text=tweet_text, images=images)
    return [chunk]
